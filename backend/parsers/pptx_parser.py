import base64
import colorsys
from collections.abc import Callable
from io import BytesIO
import mimetypes
from pathlib import Path
import shutil
import subprocess
import tempfile
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from typing import Any, Optional
import matplotlib
matplotlib.use('Agg')  # non-interactive backend, must be set before importing pyplot
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np


# ── Colour helpers ────────────────────────────────────────────────────────────

# Preset color name (a:prstClr) → hex.  Covers the most common DrawingML names.
_PRESET_COLORS: dict[str, str] = {
    'black': '#000000', 'white': '#FFFFFF', 'red': '#FF0000',
    'lime': '#00FF00', 'blue': '#0000FF', 'yellow': '#FFFF00',
    'cyan': '#00FFFF', 'magenta': '#FF00FF', 'silver': '#C0C0C0',
    'gray': '#808080', 'grey': '#808080', 'maroon': '#800000',
    'olive': '#808000', 'green': '#008000', 'purple': '#800080',
    'teal': '#008080', 'navy': '#000080', 'orange': '#FFA500',
    'darkRed': '#8B0000', 'darkGreen': '#006400', 'darkBlue': '#00008B',
    'darkCyan': '#008B8B', 'darkMagenta': '#8B008B', 'darkOrange': '#FF8C00',
    'darkGray': '#A9A9A9', 'darkGrey': '#A9A9A9', 'lightGray': '#D3D3D3',
    'lightGrey': '#D3D3D3', 'lightBlue': '#ADD8E6', 'lightGreen': '#90EE90',
    'lightYellow': '#FFFFE0', 'lightPink': '#FFB6C1', 'pink': '#FFC0CB',
    'gold': '#FFD700', 'violet': '#EE82EE', 'indigo': '#4B0082',
    'brown': '#A52A2A', 'coral': '#FF7F50', 'salmon': '#FA8072',
    'khaki': '#F0E68C', 'ivory': '#FFFFF0', 'snow': '#FFFAFA',
    'wheat': '#F5DEB3', 'beige': '#F5F5DC', 'linen': '#FAF0E6',
    'tan': '#D2B48C', 'peru': '#CD853F', 'chocolate': '#D2691E',
    'sienna': '#A0522D', 'crimson': '#DC143C', 'tomato': '#FF6347',
    'turquoise': '#40E0D0', 'aqua': '#00FFFF', 'aquaMarine': '#7FFFD4',
    'skyBlue': '#87CEEB', 'royalBlue': '#4169E1', 'steelBlue': '#4682B4',
    'slateBlue': '#6A5ACD', 'orchid': '#DA70D6', 'plum': '#DDA0DD',
    'lavender': '#E6E6FA', 'thistle': '#D8BFD8', 'mistyRose': '#FFE4E1',
}

# Scheme name aliases (OOXML spec)
_SCHEME_ALIASES: dict[str, str] = {
    'bg1': 'lt1', 'bg2': 'lt2', 'tx1': 'dk1', 'tx2': 'dk2',
}

# Module-level theme color cache — set at the start of each parse_pptx call.
# Maps scheme name (e.g. "dk1", "accent1", "bg1") → "#RRGGBB".
_theme_colors: dict[str, str] = {}


def _extract_theme_colors(prs) -> dict[str, str]:
    """Read the presentation theme and return a scheme-name → '#RRGGBB' map."""
    result: dict[str, str] = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme'
        )
        from lxml import etree as _etree
        root = _etree.fromstring(theme_part.blob)
        clr_scheme = root.find('.//' + qn('a:clrScheme'))
        if clr_scheme is not None:
            for child in clr_scheme:
                name = child.tag.split('}')[-1]
                for color_el in child:
                    ctag = color_el.tag.split('}')[-1]
                    if ctag == 'srgbClr':
                        val = color_el.get('val', '')
                        if val:
                            result[name] = f'#{val.upper()}'
                    elif ctag == 'sysClr':
                        last = color_el.get('lastClr', '')
                        if last:
                            result[name] = f'#{last.upper()}'
    except Exception:
        pass
    # Resolve aliases
    for alias, target in _SCHEME_ALIASES.items():
        if target in result and alias not in result:
            result[alias] = result[target]
    return result


def _apply_clr_mods(hex_color: str, clr_el) -> str:
    """Apply lumMod / lumOff / shade / tint child elements to *hex_color*."""
    h_str = hex_color.lstrip('#')
    r = int(h_str[0:2], 16) / 255
    g = int(h_str[2:4], 16) / 255
    b = int(h_str[4:6], 16) / 255
    h, lum, s = colorsys.rgb_to_hls(r, g, b)
    alpha = 1.0

    for child in clr_el:
        tag = child.tag.split('}')[-1]
        try:
            val = int(child.get('val', '100000')) / 100000
            if tag == 'lumMod':
                lum = lum * val
            elif tag == 'lumOff':
                lum = lum + val
            elif tag == 'shade':
                lum = lum * val
            elif tag == 'tint':
                lum = lum + (1 - lum) * val
            elif tag == 'alpha':
                alpha = val
            elif tag == 'alphaMod':
                alpha = alpha * val
            elif tag == 'alphaOff':
                alpha = alpha + val
        except Exception:
            pass

    lum = max(0.0, min(1.0, lum))
    alpha = max(0.0, min(1.0, alpha))
    r2, g2, b2 = colorsys.hls_to_rgb(h, lum, s)
    r_i = int(r2 * 255)
    g_i = int(g2 * 255)
    b_i = int(b2 * 255)
    if alpha >= 0.999:
        return f'#{r_i:02X}{g_i:02X}{b_i:02X}'
    return f'rgba({r_i}, {g_i}, {b_i}, {alpha:.3f})'


def _resolve_clr_element(clr_el) -> Optional[str]:
    """Resolve any DrawingML color element (srgbClr / schemeClr / prstClr / sysClr) to '#RRGGBB'."""
    if clr_el is None:
        return None
    tag = clr_el.tag.split('}')[-1]

    if tag == 'srgbClr':
        val = clr_el.get('val', '')
        return _apply_clr_mods(f'#{val.upper()}', clr_el) if val else None

    elif tag == 'schemeClr':
        scheme = clr_el.get('val', '')
        base = _theme_colors.get(scheme)
        if not base:
            return None
        return _apply_clr_mods(base, clr_el)

    elif tag == 'prstClr':
        name = clr_el.get('val', '')
        base = _PRESET_COLORS.get(name)
        return _apply_clr_mods(base, clr_el) if base else None

    elif tag == 'sysClr':
        last = clr_el.get('lastClr', '')
        return _apply_clr_mods(f'#{last.upper()}', clr_el) if last else None

    return None


def _solid_fill_hex(element) -> Optional[str]:
    """Return '#RRGGBB' for the first <a:solidFill> child of *element*, or None."""
    if element is None:
        return None
    solid = element.find(qn('a:solidFill'))
    if solid is None:
        return None
    for child in solid:
        result = _resolve_clr_element(child)
        if result:
            return result
    return None


def _color_hex(color_obj) -> Optional[str]:
    """Legacy helper: resolve a python-pptx ColorFormat object to '#RRGGBB'."""
    try:
        if color_obj.type is not None:
            return f"#{color_obj.rgb}"
    except Exception:
        pass
    # Fallback: parse the underlying XML element
    try:
        el = color_obj._color  # the raw color element (e.g. <a:schemeClr>)
        return _resolve_clr_element(el)
    except Exception:
        pass
    return None


def _fill_hex(fill_obj) -> Optional[str]:
    # Try python-pptx API first (works for plain RGB fills)
    try:
        from pptx.enum.dml import MSO_FILL
        if fill_obj.type == MSO_FILL.SOLID:
            result = _color_hex(fill_obj.fore_color)
            if result:
                return result
    except Exception:
        pass
    # Fallback: parse the XML element directly (handles schemeClr / prstClr)
    try:
        return _solid_fill_hex(fill_obj._fill)
    except Exception:
        pass
    return None


def _shape_fill_hex(shape) -> Optional[str]:
    """Extract fill color from a shape, using XML as primary source."""
    try:
        # Direct spPr solidFill
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is None:
            # Some shape types nest spPr differently
            spPr = shape._element.find('.//' + qn('p:spPr'))
        if spPr is not None:
            result = _solid_fill_hex(spPr)
            if result:
                return result
    except Exception:
        pass
    # Fallback via python-pptx fill API
    try:
        return _fill_hex(shape.fill)
    except Exception:
        pass
    return None


def _gradient_fill_css(fill_parent) -> Optional[str]:
    if fill_parent is None:
        return None

    grad = fill_parent.find(qn('a:gradFill'))
    if grad is None:
        return None

    stops_parent = grad.find(qn('a:gsLst'))
    if stops_parent is None:
        return None

    stops: list[tuple[float, str]] = []
    for stop in stops_parent:
        try:
            pos = int(stop.get('pos', '0')) / 1000
        except Exception:
            pos = 0.0
        color = None
        for child in stop:
            color = _resolve_clr_element(child)
            if color:
                break
        if color:
            stops.append((max(0.0, min(100.0, pos)), color))

    if not stops:
        return None

    # Degenerate gradients are visually just a solid fill.
    if len({color for _, color in stops}) == 1:
        return stops[0][1]

    linear = grad.find(qn('a:lin'))
    if linear is not None:
        try:
            angle = int(linear.get('ang', '0')) / 60000
        except Exception:
            angle = 0
        stop_list = ", ".join(f"{color} {pos:.0f}%" for pos, color in stops)
        return f"linear-gradient({angle:.2f}deg, {stop_list})"

    path = grad.find(qn('a:path'))
    if path is not None:
        stop_list = ", ".join(f"{color} {pos:.0f}%" for pos, color in stops)
        return f"radial-gradient(circle, {stop_list})"

    stop_list = ", ".join(f"{color} {pos:.0f}%" for pos, color in stops)
    return f"linear-gradient(0deg, {stop_list})"


def _shape_fill_gradient(shape) -> Optional[str]:
    try:
        spPr = _shape_sp_pr(shape)
        if spPr is not None:
            return _gradient_fill_css(spPr)
    except Exception:
        pass
    return None


def _shape_line_hex(shape) -> Optional[str]:
    try:
        return _color_hex(shape.line.color)
    except Exception:
        pass
    try:
        ln = shape._element.find('.//' + qn('a:ln'))
        if ln is not None:
            return _solid_fill_hex(ln)
    except Exception:
        pass
    return None


def _shape_line_width(shape) -> Optional[int]:
    try:
        width = shape.line.width
        return int(width) if width is not None else None
    except Exception:
        pass
    try:
        ln = shape._element.find('.//' + qn('a:ln'))
        if ln is not None and ln.get('w'):
            return int(ln.get('w'))
    except Exception:
        pass
    return None


def _shape_sp_pr(shape):
    try:
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = shape._element.find('.//' + qn('p:spPr'))
        return spPr
    except Exception:
        return None


def _shape_transform(shape) -> dict[str, Any]:
    rotation = 0.0
    flip_horizontal = False
    flip_vertical = False
    try:
        spPr = _shape_sp_pr(shape)
        xfrm = None if spPr is None else spPr.find(qn('a:xfrm'))
        if xfrm is not None:
            if xfrm.get('rot'):
                rotation = int(xfrm.get('rot', '0')) / 60000
            flip_horizontal = xfrm.get('flipH') == '1'
            flip_vertical = xfrm.get('flipV') == '1'
    except Exception:
        pass
    return {
        "rotation": rotation,
        "flip_horizontal": flip_horizontal,
        "flip_vertical": flip_vertical,
    }


def _shape_preset_geometry(shape) -> Optional[str]:
    try:
        spPr = _shape_sp_pr(shape)
        geom = None if spPr is None else spPr.find(qn('a:prstGeom'))
        if geom is not None:
            return geom.get('prst')
    except Exception:
        pass
    return None


def _geom_env(width: int, height: int) -> dict[str, int]:
    env = {
        'w': width,
        'h': height,
        'l': 0,
        't': 0,
        'r': width,
        'b': height,
        'wd2': width // 2,
        'wd4': width // 4,
        'hd2': height // 2,
        'hd4': height // 4,
        'hc': width // 2,
        'vc': height // 2,
        'ss': min(width, height),
        'ls': max(width, height),
    }
    return env


def _geom_value(token: str, env: dict[str, int]) -> int:
    try:
        return int(token)
    except Exception:
        return env.get(token, 0)


def _eval_geom_formula(formula: str, env: dict[str, int]) -> int:
    parts = formula.split()
    if not parts:
        return 0
    op = parts[0]
    args = parts[1:]

    try:
        if op == 'val' and len(args) >= 1:
            return _geom_value(args[0], env)
        if op == '*/' and len(args) >= 3:
            a, b, c = (_geom_value(arg, env) for arg in args[:3])
            return int(round((a * b) / c)) if c else 0
        if op == '+-' and len(args) >= 3:
            a, b, c = (_geom_value(arg, env) for arg in args[:3])
            return a + b - c
        if op == '+/' and len(args) >= 3:
            a, b, c = (_geom_value(arg, env) for arg in args[:3])
            return int(round((a + b) / c)) if c else 0
        if op == 'max' and len(args) >= 2:
            return max(_geom_value(args[0], env), _geom_value(args[1], env))
        if op == 'min' and len(args) >= 2:
            return min(_geom_value(args[0], env), _geom_value(args[1], env))
    except Exception:
        return 0

    return _geom_value(args[0], env) if args else 0


def _shape_geometry_adjustments(shape) -> Optional[dict[str, int]]:
    try:
        spPr = _shape_sp_pr(shape)
        prst_geom = None if spPr is None else spPr.find(qn('a:prstGeom'))
        if prst_geom is None:
            return None
        avLst = prst_geom.find(qn('a:avLst'))
        if avLst is None:
            return None
        env = _geom_env(100000, 100000)
        adjustments: dict[str, int] = {}
        for gd in avLst:
            name = gd.get('name')
            formula = gd.get('fmla', '')
            if not name or not formula:
                continue
            value = _eval_geom_formula(formula, env)
            env[name] = value
            adjustments[name] = value
        return adjustments or None
    except Exception:
        return None


def _shape_svg_path(shape) -> tuple[Optional[str], Optional[int], Optional[int]]:
    try:
        spPr = _shape_sp_pr(shape)
        cust_geom = None if spPr is None else spPr.find(qn('a:custGeom'))
        if cust_geom is None:
            return None, None, None

        shape_width = int(shape.width or 0) or 1
        shape_height = int(shape.height or 0) or 1

        gdLst = cust_geom.find(qn('a:gdLst'))
        path_lst = cust_geom.find(qn('a:pathLst'))
        if path_lst is None:
            return None, None, None

        path_segments: list[str] = []
        for path in path_lst:
            path_width = int(path.get('w', str(shape_width)) or shape_width)
            path_height = int(path.get('h', str(shape_height)) or shape_height)
            env = _geom_env(path_width, path_height)

            if gdLst is not None:
                for gd in gdLst:
                    name = gd.get('name')
                    formula = gd.get('fmla', '')
                    if not name or not formula:
                        continue
                    env[name] = _eval_geom_formula(formula, env)

            def scale_x(raw: str) -> float:
                return (_geom_value(raw, env) * shape_width) / path_width if path_width else 0.0

            def scale_y(raw: str) -> float:
                return (_geom_value(raw, env) * shape_height) / path_height if path_height else 0.0

            for cmd in path:
                tag = cmd.tag.split('}')[-1]
                if tag in {'moveTo', 'lnTo'}:
                    pt = cmd.find(qn('a:pt'))
                    if pt is None:
                        continue
                    x = scale_x(pt.get('x', '0'))
                    y = scale_y(pt.get('y', '0'))
                    path_segments.append(f"{'M' if tag == 'moveTo' else 'L'} {x:.2f} {y:.2f}")
                elif tag == 'cubicBezTo':
                    pts = cmd.findall(qn('a:pt'))
                    if len(pts) != 3:
                        continue
                    p = [(scale_x(pt.get('x', '0')), scale_y(pt.get('y', '0'))) for pt in pts]
                    path_segments.append(
                        f"C {p[0][0]:.2f} {p[0][1]:.2f} {p[1][0]:.2f} {p[1][1]:.2f} {p[2][0]:.2f} {p[2][1]:.2f}"
                    )
                elif tag == 'close':
                    path_segments.append('Z')

        if not path_segments:
            return None, None, None

        return " ".join(path_segments), shape_width, shape_height
    except Exception:
        return None, None, None


# ── Image extraction ──────────────────────────────────────────────────────────

_BROWSER_SAFE_IMAGE_TYPES = {
    "image/apng",
    "image/avif",
    "image/bmp",
    "image/gif",
    "image/jpeg",
    "image/jpg",
    "image/png",
    "image/svg+xml",
    "image/webp",
}

_IMAGE_SUFFIX_MAP = {
    "application/octet-stream": ".bin",
    "application/x-emf": ".emf",
    "application/x-wmf": ".wmf",
    "image/bmp": ".bmp",
    "image/emf": ".emf",
    "image/gif": ".gif",
    "image/jpeg": ".jpg",
    "image/jpg": ".jpg",
    "image/png": ".png",
    "image/svg+xml": ".svg",
    "image/tiff": ".tiff",
    "image/wmf": ".wmf",
    "image/x-emf": ".emf",
    "image/x-wmf": ".wmf",
}

_EXTERNAL_IMAGE_CONVERSION_TIMEOUT_SECONDS = 2
# Images are scaled down to this dimension on the longest side before encoding.
# Slides are typically ≤1280px wide in the browser so anything larger is wasted bandwidth.
_MAX_IMAGE_DIM = 1280


def _normalize_content_type(content_type: Optional[str]) -> Optional[str]:
    if not content_type:
        return None
    return content_type.split(";", 1)[0].strip().lower() or None


def _data_url(blob: bytes, content_type: Optional[str]) -> str:
    mime = _normalize_content_type(content_type) or "application/octet-stream"
    data = base64.b64encode(blob).decode("utf-8")
    return f"data:{mime};base64,{data}"


def _guess_image_suffix(content_type: Optional[str], filename: Optional[str] = None) -> str:
    mime = _normalize_content_type(content_type)
    if mime in _IMAGE_SUFFIX_MAP:
        return _IMAGE_SUFFIX_MAP[mime]

    if filename:
        suffix = Path(filename).suffix.lower()
        if suffix:
            return suffix

    guessed = mimetypes.guess_extension(mime or "")
    if guessed == ".jpe":
        return ".jpg"
    return guessed or ".bin"


def _pillow_normalize(blob: bytes) -> Optional[tuple[bytes, str]]:
    """Open *blob* with Pillow, resize to _MAX_IMAGE_DIM, and return (bytes, mime).

    Returns JPEG at 85 % quality for opaque images (much smaller than PNG),
    or PNG for images with transparency.  Returns None if Pillow cannot open it.
    """
    try:
        with Image.open(BytesIO(blob)) as image:
            image.load()
            needs_alpha = image.mode in {"RGBA", "LA", "PA"} or ("transparency" in image.info)

            # Downscale if either dimension exceeds the cap
            w, h = image.size
            if max(w, h) > _MAX_IMAGE_DIM:
                scale = _MAX_IMAGE_DIM / max(w, h)
                image = image.resize(
                    (max(1, int(w * scale)), max(1, int(h * scale))),
                    Image.LANCZOS,
                )

            output = BytesIO()
            if needs_alpha:
                image.convert("RGBA").save(output, format="PNG", optimize=True)
                return output.getvalue(), "image/png"
            else:
                image.convert("RGB").save(output, format="JPEG", quality=85, optimize=True)
                return output.getvalue(), "image/jpeg"
    except Exception:
        return None


def _convert_with_magick(
    blob: bytes,
    content_type: Optional[str],
    filename: Optional[str] = None,
) -> Optional[bytes]:
    magick = shutil.which("magick")
    if not magick:
        return None

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            src = Path(tmpdir) / f"source{_guess_image_suffix(content_type, filename)}"
            dst = Path(tmpdir) / "converted.png"
            src.write_bytes(blob)
            subprocess.run(
                [magick, str(src), str(dst)],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=_EXTERNAL_IMAGE_CONVERSION_TIMEOUT_SECONDS,
            )
            if dst.exists():
                return dst.read_bytes()
    except Exception:
        return None
    return None


def _image_src_from_blob(
    blob: bytes,
    content_type: Optional[str],
    filename: Optional[str] = None,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
) -> str:
    # Always normalise through Pillow: resizes to _MAX_IMAGE_DIM and converts to
    # JPEG/PNG.  This is fast, keeps the response small, and handles TIFF/BMP too.
    result = _pillow_normalize(blob)
    if result:
        out_bytes, out_mime = result
        if asset_resolver:
            return asset_resolver(out_bytes, out_mime, filename)
        return _data_url(out_bytes, out_mime)

    # Pillow failed (e.g. EMF/WMF vector formats). Try ImageMagick if available.
    mime = _normalize_content_type(content_type)
    converted = _convert_with_magick(blob, mime, filename)
    if converted:
        # The magick output is a PNG — normalise it too so it gets resized.
        result2 = _pillow_normalize(converted)
        if result2:
            out_bytes, out_mime = result2
            if asset_resolver:
                return asset_resolver(out_bytes, out_mime, filename)
            return _data_url(out_bytes, out_mime)
        if asset_resolver:
            return asset_resolver(converted, "image/png", filename)
        return _data_url(converted, "image/png")

    # Could not convert at all — pass raw bytes through as a last resort.
    if asset_resolver:
        return asset_resolver(blob, mime, filename)
    return _data_url(blob, mime)

def _blip_image_src(
    elem,
    part,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
) -> Optional[str]:
    """
    Search *elem* (any lxml element) for an <a:blip r:embed="…"> and return
    a base64 data-URL for the referenced image, or None.
    """
    try:
        blip = elem.find('.//' + qn('a:blip'))
        if blip is None:
            return None
        rel_id = blip.get(qn('r:embed')) or blip.get(qn('r:link'))
        if not rel_id:
            return None
        rel = part.rels.get(rel_id)
        if rel is not None and getattr(rel, "is_external", False):
            target_ref = getattr(rel, "target_ref", None)
            if isinstance(target_ref, str) and target_ref.startswith(("http://", "https://", "data:")):
                return target_ref
            return None

        img_part = part.related_part(rel_id)
        partname = getattr(img_part, "partname", None)
        filename = Path(str(partname)).name if partname is not None else None
        return _image_src_from_blob(
            img_part.blob,
            getattr(img_part, "content_type", None),
            filename,
            asset_resolver,
        )
    except Exception:
        return None


def _scaled_emu(value: int, scale: float) -> int:
    return int(round(value * scale))


def _geometry(
    shape,
    left_offset: float,
    top_offset: float,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    child_origin_x: int = 0,
    child_origin_y: int = 0,
) -> dict:
    geometry = {
        "left":   _scaled_emu((shape.left  or 0) - child_origin_x, scale_x) + int(round(left_offset)),
        "top":    _scaled_emu((shape.top   or 0) - child_origin_y, scale_y) + int(round(top_offset)),
        "width":  _scaled_emu(shape.width  or 0, scale_x),
        "height": _scaled_emu(shape.height or 0, scale_y),
    }
    geometry.update(_shape_transform(shape))
    return geometry


def _image_dict(
    shape,
    idx: int,
    left_offset: float,
    top_offset: float,
    src: str,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    child_origin_x: int = 0,
    child_origin_y: int = 0,
) -> dict:
    return {
        "index": idx,
        "name": shape.name,
        "shape_type": "image",
        "text": "",
        "image_src": src,
        "paragraphs": [],
        "fill_color": None,
        "ph_idx": None,
        "vertical_anchor": "top",
        **_geometry(
            shape,
            left_offset,
            top_offset,
            scale_x,
            scale_y,
            child_origin_x,
            child_origin_y,
        ),
    }


def _decoration_dict(
    shape,
    idx: int,
    left_offset: float,
    top_offset: float,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    child_origin_x: int = 0,
    child_origin_y: int = 0,
) -> Optional[dict]:
    fill_color = _shape_fill_hex(shape)
    fill_gradient = _shape_fill_gradient(shape)
    if not fill_color and fill_gradient and 'gradient(' not in fill_gradient:
        fill_color = fill_gradient
        fill_gradient = None
    stroke_color = _shape_line_hex(shape)
    stroke_width = _shape_line_width(shape)
    preset_geometry = _shape_preset_geometry(shape)
    svg_path, svg_viewbox_width, svg_viewbox_height = _shape_svg_path(shape)
    geometry_adjustments = _shape_geometry_adjustments(shape)
    if not fill_color and not fill_gradient and not stroke_color:
        return None

    return {
        "index": idx,
        "name": shape.name,
        "shape_type": "decoration",
        "text": "",
        "image_src": None,
        "paragraphs": [],
        "fill_color": fill_color,
        "fill_gradient": fill_gradient,
        "stroke_color": stroke_color,
        "stroke_width": stroke_width,
        "preset_geometry": preset_geometry,
        "geometry_adjustments": geometry_adjustments,
        "svg_path": svg_path,
        "svg_viewbox_width": svg_viewbox_width,
        "svg_viewbox_height": svg_viewbox_height,
        "ph_idx": None,
        "vertical_anchor": "top",
        **_geometry(
            shape,
            left_offset,
            top_offset,
            scale_x,
            scale_y,
            child_origin_x,
            child_origin_y,
        ),
    }


# ── Text extraction ───────────────────────────────────────────────────────────

_ALIGN_MAP = {
    PP_ALIGN.CENTER:    "center",
    PP_ALIGN.RIGHT:     "right",
    PP_ALIGN.JUSTIFY:   "justify",
    PP_ALIGN.DISTRIBUTE:"justify",
}

_ANCHOR_MAP = {
    MSO_ANCHOR.MIDDLE: "middle",
    MSO_ANCHOR.BOTTOM: "bottom",
}


def _paragraph_has_bullet(para) -> bool:
    try:
        pPr = para._p.pPr
        if pPr is None:
            return False
        for child in pPr:
            tag = getattr(child, "tag", "")
            if tag in (qn('a:buChar'), qn('a:buAutoNum')):
                return True
    except Exception:
        pass
    return False


def _text_dict(shape, idx: int) -> Optional[dict]:
    try:
        has_tf = shape.has_text_frame
    except Exception:
        return None
    if not has_tf:
        return None

    try:
        tf = shape.text_frame
    except Exception:
        return None

    paragraphs = []

    for para in tf.paragraphs:
        runs = []
        for run in para.runs:
            try:
                f = run.font
                runs.append({
                    "text":      run.text,
                    "bold":      f.bold,
                    "italic":    f.italic,
                    "underline": f.underline,
                    "strike":    run._r.get('strike') not in (None, 'noStrike'),
                    "font_name": f.name,
                    "size":      f.size.pt if f.size is not None else None,
                    "color":     _solid_fill_hex(run._r.find(qn('a:rPr'))) or _color_hex(f.color),
                })
            except Exception:
                try:
                    runs.append({"text": run.text, "bold": None, "italic": None,
                                 "underline": None, "strike": None, "font_name": None, "size": None, "color": None})
                except Exception:
                    pass

        try:
            para_text = para.text
        except Exception:
            continue

        if not runs and para_text:
            runs = [{"text": para_text, "bold": None, "italic": None,
                     "underline": None, "strike": None, "font_name": None, "size": None, "color": None}]

        try:
            align = _ALIGN_MAP.get(para.alignment, "left")
        except Exception:
            align = "left"

        paragraphs.append({
            "text":  para_text,
            "align": align,
            "runs":  runs,
            "bullet": _paragraph_has_bullet(para),
            "level": getattr(para, "level", 0) or 0,
        })

    full_text = "\n".join(p["text"] for p in paragraphs)
    if not full_text.strip():
        return None

    ph_idx = None
    try:
        if shape.is_placeholder:
            ph_idx = shape.placeholder_format.idx
    except Exception:
        pass

    vertical_anchor = "top"
    try:
        vertical_anchor = _ANCHOR_MAP.get(tf.vertical_anchor, "top")
    except Exception:
        pass

    return {
        "index": idx,
        "name": shape.name,
        "shape_type": "text",
        "text": full_text,
        "paragraphs": paragraphs,
        "fill_color": _shape_fill_hex(shape),
        "fill_gradient": _shape_fill_gradient(shape),
        "stroke_color": _shape_line_hex(shape),
        "stroke_width": _shape_line_width(shape),
        "ph_idx": ph_idx,
        "vertical_anchor": vertical_anchor,
        "image_src": None,
    }


# ── Shape dispatcher ──────────────────────────────────────────────────────────

_CHART_COLORS = ['#4472C4', '#ED7D31', '#A9D18E', '#FFC000', '#5B9BD5', '#70AD47', '#FF0000', '#00B0F0']


def _render_chart_png(shape) -> Optional[bytes]:
    """Render a pptx chart shape to a PNG using matplotlib. Returns PNG bytes or None."""
    try:
        chart = shape.chart
        from pptx.enum.chart import XL_CHART_TYPE

        series_list = list(chart.series)
        if not series_list:
            return None

        # Collect categories from first series that has them
        categories: list[str] = []
        try:
            plot = chart.plots[0]
            cats = plot.series[0].data_labels  # may not exist
        except Exception:
            pass
        try:
            from pptx.oxml.ns import qn as _qn
            ser0 = series_list[0]
            cat_el = ser0._element.find('.//' + _qn('c:cat'))
            if cat_el is not None:
                for pt in cat_el.iter(_qn('c:pt')):
                    v = pt.find(_qn('c:v'))
                    categories.append(v.text if v is not None else '')
        except Exception:
            pass

        x = np.arange(max(len(list(s.values)) for s in series_list))
        if not categories:
            categories = [str(i + 1) for i in range(len(x))]

        chart_type_val = chart.chart_type
        is_horizontal = chart_type_val in (
            getattr(XL_CHART_TYPE, 'BAR_CLUSTERED', None),
            getattr(XL_CHART_TYPE, 'BAR_STACKED', None),
            getattr(XL_CHART_TYPE, 'BAR_STACKED_100', None),
        )
        is_stacked = 'STACKED' in str(chart_type_val)
        is_line = 'LINE' in str(chart_type_val)
        is_pie = 'PIE' in str(chart_type_val) or 'DOUGHNUT' in str(chart_type_val)

        fig, ax = plt.subplots(figsize=(8, 4.5), dpi=120)
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#f8f9fa')
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        bottom = np.zeros(len(x))

        for i, series in enumerate(series_list):
            color = _CHART_COLORS[i % len(_CHART_COLORS)]
            values = [v if v is not None else 0 for v in series.values]
            # Pad to x length
            values = values[:len(x)]
            values += [0] * (len(x) - len(values))
            vals = np.array(values, dtype=float)
            label = series.name or f'Series {i + 1}'

            if is_pie:
                ax.pie(np.abs(vals), labels=categories[:len(vals)], autopct='%1.0f%%',
                       colors=_CHART_COLORS[:len(vals)])
                break
            elif is_line:
                ax.plot(x, vals, color=color, marker='o', linewidth=2, markersize=4, label=label)
            elif is_horizontal:
                if is_stacked:
                    ax.barh(x, vals, left=bottom, color=color, label=label, height=0.6)
                    bottom += vals
                else:
                    w = 0.8 / len(series_list)
                    ax.barh(x + i * w - 0.4 + w / 2, vals, color=color, label=label, height=w)
            else:
                if is_stacked:
                    ax.bar(x, vals, bottom=bottom, color=color, label=label, width=0.6)
                    bottom += vals
                else:
                    w = 0.8 / len(series_list)
                    ax.bar(x + i * w - 0.4 + w / 2, vals, color=color, label=label, width=w)

        if not is_pie:
            if is_horizontal:
                ax.set_yticks(x)
                ax.set_yticklabels(categories[:len(x)], fontsize=8)
            else:
                ax.set_xticks(x)
                ax.set_xticklabels(categories[:len(x)], fontsize=8, rotation=30, ha='right')
            if len(series_list) > 1:
                ax.legend(fontsize=8, framealpha=0.7)

        try:
            title = chart.chart_title.text_frame.text
            if title:
                ax.set_title(title, fontsize=10, fontweight='bold')
        except Exception:
            pass

        fig.tight_layout(pad=1.0)
        buf = BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', facecolor='white')
        plt.close(fig)
        return buf.getvalue()
    except Exception:
        plt.close('all')
        return None


def _parse_table_shape(shape, idx: int) -> Optional[dict]:
    """Extract a pptx TABLE shape into a structured dict."""
    try:
        tbl = shape.table
        rows_data = []
        text_lines = []
        for row in tbl.rows:
            cells = []
            row_texts = []
            for cell in row.cells:
                try:
                    text = cell.text or ''
                except Exception:
                    text = ''
                row_texts.append(text.strip())
                try:
                    # Try XML-based extraction first (handles scheme/preset colors)
                    tcPr = cell._tc.find(qn('a:tcPr'))
                    fill = _solid_fill_hex(tcPr) or _fill_hex(cell.fill)
                except Exception:
                    fill = None
                cells.append({'text': text, 'fill': fill})
            rows_data.append(cells)
            if any(row_texts):
                text_lines.append(" | ".join(row_texts))
        return {
            'index': idx,
            'name': shape.name,
            'shape_type': 'table',
            'text': "\n".join(text_lines),
            'paragraphs': [],
            'fill_color': None,
            'ph_idx': None,
            'vertical_anchor': 'top',
            'image_src': None,
            'table_data': rows_data,
        }
    except Exception:
        return None


def _parse_shape(
    shape,
    idx: int,
    part,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
    left_offset: float = 0,
    top_offset: float = 0,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    child_origin_x: int = 0,
    child_origin_y: int = 0,
) -> list[dict]:
    try:
        return _parse_shape_inner(shape, idx, part, asset_resolver, left_offset, top_offset, scale_x, scale_y, child_origin_x, child_origin_y)
    except Exception:
        return []


def _parse_shape_inner(
    shape,
    idx: int,
    part,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
    left_offset: float = 0,
    top_offset: float = 0,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    child_origin_x: int = 0,
    child_origin_y: int = 0,
) -> list[dict]:
    stype = shape.shape_type

    # ── Direct picture shape ──────────────────────────────────────────────────
    if stype == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            src = _image_src_from_blob(
                img.blob,
                getattr(img, "content_type", None),
                getattr(img, "filename", None),
                asset_resolver,
            )
            return [
                _image_dict(
                    shape,
                    idx,
                    left_offset,
                    top_offset,
                    src,
                    scale_x,
                    scale_y,
                    child_origin_x,
                    child_origin_y,
                )
            ]
        except Exception:
            pass
        # Fallback: try blip extraction
        src = _blip_image_src(shape._element, part, asset_resolver)
        if src:
            return [
                _image_dict(
                    shape,
                    idx,
                    left_offset,
                    top_offset,
                    src,
                    scale_x,
                    scale_y,
                    child_origin_x,
                    child_origin_y,
                )
            ]
        return []

    # ── Group: recurse with offset ─────────────────────────────────────────────
    if stype == MSO_SHAPE_TYPE.GROUP:
        # Group children use the group's own child coordinate space, which can be
        # translated and scaled relative to the slide.
        gl = _scaled_emu((shape.left or 0) - child_origin_x, scale_x) + left_offset
        gt = _scaled_emu((shape.top  or 0) - child_origin_y, scale_y) + top_offset
        group_origin_x = getattr(getattr(shape._element, "chOff", None), "x", 0) or 0
        group_origin_y = getattr(getattr(shape._element, "chOff", None), "y", 0) or 0
        group_extent_x = getattr(getattr(shape._element, "chExt", None), "cx", 0) or 0
        group_extent_y = getattr(getattr(shape._element, "chExt", None), "cy", 0) or 0
        group_scale_x = scale_x * ((shape.width or 0) / group_extent_x) if group_extent_x else scale_x
        group_scale_y = scale_y * ((shape.height or 0) / group_extent_y) if group_extent_y else scale_y
        children = []
        for j, child in enumerate(shape.shapes):
            children.extend(
                _parse_shape(
                    child,
                    idx * 10000 + j,
                    part,
                    asset_resolver,
                    gl,
                    gt,
                    group_scale_x,
                    group_scale_y,
                    group_origin_x,
                    group_origin_y,
                )
            )
        return children

    # ── Chart ─────────────────────────────────────────────────────────────────
    if getattr(shape, 'has_chart', False):
        png = _render_chart_png(shape)
        if png:
            # Normalise (resize) the rendered PNG before encoding
            result2 = _pillow_normalize(png)
            if result2:
                png, _ = result2
            src = _data_url(png, 'image/png')
        else:
            # Fallback: grey placeholder
            src = 'data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=='
        return [_image_dict(shape, idx, left_offset, top_offset, src, scale_x, scale_y, child_origin_x, child_origin_y)]

    # ── PPTX table ────────────────────────────────────────────────────────────
    if stype == MSO_SHAPE_TYPE.TABLE:
        result = _parse_table_shape(shape, idx)
        if result:
            geo = _geometry(shape, left_offset, top_offset, scale_x, scale_y, child_origin_x, child_origin_y)
            result.update(geo)
            return [result]
        return []

    decoration = _decoration_dict(shape, idx, left_offset, top_offset, scale_x, scale_y, child_origin_x, child_origin_y)
    if decoration:
        return [decoration]

    # ── Any shape with a blip (picture) fill ──────────────────────────────────
    # Covers picture placeholders, rectangles with image fill, etc.
    src = _blip_image_src(shape._element, part, asset_resolver)
    if src:
        return [
            _image_dict(
                shape,
                idx,
                left_offset,
                top_offset,
                src,
                scale_x,
                scale_y,
                child_origin_x,
                child_origin_y,
            )
        ]

    # ── Text shape ────────────────────────────────────────────────────────────
    result = _text_dict(shape, idx)
    if result:
        geo = _geometry(
            shape,
            left_offset,
            top_offset,
            scale_x,
            scale_y,
            child_origin_x,
            child_origin_y,
        )
        result.update(geo)
        return [result]

    return []


def _parse_shape_collection(
    shape_collection,
    part,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
    idx_offset: int = 0,
    images_only: bool = False,
) -> list[dict]:
    parsed = []
    for j, shape in enumerate(shape_collection):
        entries = _parse_shape(shape, idx_offset + j, part, asset_resolver)
        if images_only:
            entries = [entry for entry in entries if entry["shape_type"] == "image"]
        parsed.extend(entries)
    return parsed


def _explicit_bg_pr(base_slide) -> Any | None:
    try:
        bg = base_slide._element.cSld.bg
    except Exception:
        return None
    return None if bg is None else bg.bgPr


def _background_sources(slide) -> list[tuple[Any, Any]]:
    sources = [(slide, slide.part)]
    try:
        layout = slide.slide_layout
        sources.append((layout, layout.part))
        master = layout.slide_master
        sources.append((master, master.part))
    except Exception:
        pass
    return sources


def _parse_inherited_visuals(
    slide,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
) -> list[dict]:
    visuals: list[dict] = []
    seen: set[tuple[Any, ...]] = set()

    try:
        inherited_sources = [
            (slide.slide_layout.slide_master.shapes, slide.slide_layout.slide_master.part, 1_000_000),
            (slide.slide_layout.shapes, slide.slide_layout.part, 2_000_000),
        ]
    except Exception:
        return visuals

    for shape_collection, part, idx_offset in inherited_sources:
        try:
            parsed = _parse_shape_collection(shape_collection, part, asset_resolver, idx_offset=idx_offset)
        except Exception:
            continue

        for shape in parsed:
            if shape["shape_type"] not in {"image", "decoration"}:
                continue

            try:
                original_idx = shape["index"] - idx_offset
                source_shape = shape_collection[original_idx]
                if getattr(source_shape, "is_placeholder", False):
                    continue
            except Exception:
                pass

            key = (
                shape["shape_type"],
                shape.get("name"),
                shape.get("left"),
                shape.get("top"),
                shape.get("width"),
                shape.get("height"),
                shape.get("image_src"),
                shape.get("fill_color"),
                shape.get("stroke_color"),
            )
            if key in seen:
                continue
            seen.add(key)
            visuals.append(shape)

    return visuals


def _looks_like_slide_background(shape: dict, slide_width: int, slide_height: int) -> bool:
    if shape.get("shape_type") != "image":
        return False

    width = shape.get("width", 0)
    height = shape.get("height", 0)
    left = abs(shape.get("left", 0))
    top = abs(shape.get("top", 0))
    if slide_width <= 0 or slide_height <= 0:
        return False

    return (
        left <= slide_width * 0.02 and
        top <= slide_height * 0.02 and
        width >= slide_width * 0.9 and
        height >= slide_height * 0.9
    )


def _resolve_background_image_src(
    slide,
    slide_width: int,
    slide_height: int,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
) -> Optional[str]:
    for base_slide, part in _background_sources(slide):
        bg_pr = _explicit_bg_pr(base_slide)
        if bg_pr is None:
            if base_slide is not slide:
                shapes = _parse_shape_collection(base_slide.shapes, part, asset_resolver, images_only=True)
                for shape in shapes:
                    if _looks_like_slide_background(shape, slide_width, slide_height):
                        return shape["image_src"]
            continue

        src = _blip_image_src(bg_pr, part, asset_resolver)
        if src:
            return src
    return None


def _resolve_background_color(slide) -> Optional[str]:
    from pptx.dml.fill import FillFormat

    for base_slide, _ in _background_sources(slide):
        bg_pr = _explicit_bg_pr(base_slide)
        if bg_pr is None:
            continue
        try:
            fill = FillFormat.from_fill_parent(bg_pr)
            color = _fill_hex(fill)
            if color:
                return color
        except Exception:
            pass
    return None


def get_pptx_table_cells(file_path: str, slide_index: int, table_index: int) -> list[dict]:
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return []

    slide = prs.slides[slide_index]
    if table_index < 0 or table_index >= len(slide.shapes):
        return []

    shape = slide.shapes[table_index]
    if not getattr(shape, "has_table", False):
        return []

    cells: list[dict] = []
    for row_idx, row in enumerate(shape.table.rows):
        for cell_idx, cell in enumerate(row.cells):
            text = (getattr(cell, "text", "") or "").strip()
            cells.append({"row": row_idx, "col": cell_idx, "text": text})
    return cells


def get_pptx_cell_text(
    file_path: str,
    slide_index: int,
    table_index: int,
    row_index: int,
    cell_index: int,
) -> str | None:
    prs = Presentation(file_path)
    if slide_index < 0 or slide_index >= len(prs.slides):
        return None

    slide = prs.slides[slide_index]
    if table_index < 0 or table_index >= len(slide.shapes):
        return None

    shape = slide.shapes[table_index]
    if not getattr(shape, "has_table", False):
        return None

    table = shape.table
    if row_index < 0 or row_index >= len(table.rows):
        return None

    row = table.rows[row_index]
    if cell_index < 0 or cell_index >= len(row.cells):
        return None

    return getattr(row.cells[cell_index], "text", None)


# ── Public API ────────────────────────────────────────────────────────────────

def parse_pptx(
    file_path: str,
    asset_resolver: Optional[Callable[[bytes, Optional[str], Optional[str]], str]] = None,
) -> dict[str, Any]:
    global _theme_colors
    prs = Presentation(file_path)
    _theme_colors = _extract_theme_colors(prs)
    slides = []

    for i, slide in enumerate(prs.slides):
        try:
            background_image_src = _resolve_background_image_src(
                slide,
                prs.slide_width,
                prs.slide_height,
                asset_resolver,
            )
        except Exception:
            background_image_src = None

        shapes = []
        try:
            shapes.extend(_parse_inherited_visuals(slide, asset_resolver))
        except Exception:
            pass
        try:
            shapes.extend(_parse_shape_collection(slide.shapes, slide.part, asset_resolver))
        except Exception:
            pass

        try:
            bg_color = _resolve_background_color(slide)
        except Exception:
            bg_color = None

        slides.append({
            "index":                i,
            "shapes":               shapes,
            "background":           bg_color,
            "background_image_src": background_image_src,
        })

    return {
        "slides":       slides,
        "total":        len(prs.slides),
        "slide_width":  prs.slide_width,
        "slide_height": prs.slide_height,
    }
