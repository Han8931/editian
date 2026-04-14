import copy
import shutil
from typing import Any
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
from lxml import etree


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def apply_pptx_revisions(
    input_path: str,
    output_path: str,
    revisions: list[dict[str, Any]],
) -> None:
    if input_path != output_path:
        shutil.copy2(input_path, output_path)

    prs = Presentation(output_path)

    for revision in revisions:
        scope = revision["scope"]
        revised_text: str = revision.get("revised", "")

        font_name: str | None   = revision.get("font_name")
        font_size: float | None = revision.get("font_size")
        align: str | None       = revision.get("align")
        bold: bool | None       = revision.get("bold")
        italic: bool | None     = revision.get("italic")
        underline: bool | None  = revision.get("underline")
        strike: bool | None     = revision.get("strike")
        bullet: bool | None     = revision.get("bullet")

        stype = scope["type"]

        # ── Slide-level structural operations ────────────────────────────────
        if stype == "insert_slide":
            _insert_blank_slide(prs, scope.get("slide_index", len(prs.slides) - 1))
            continue

        if stype == "delete_slide":
            idx = scope.get("slide_index")
            if idx is not None:
                _delete_slide(prs, idx)
            continue

        if stype == "duplicate_slide":
            idx = scope.get("slide_index")
            if idx is not None:
                _duplicate_slide(prs, idx)
            continue

        # ── Text / formatting operations ─────────────────────────────────────
        if stype == "shape":
            slide_idx = scope.get("slide_index")
            shape_indices = scope.get("shape_indices") or []
            if slide_idx is not None and slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                for shape_idx in shape_indices:
                    if shape_idx < len(slide.shapes):
                        s = slide.shapes[shape_idx]
                        if s.has_text_frame:
                            _set_shape_text(s, revised_text, font_name, font_size, align, bold, italic, underline, strike, bullet)

        elif stype == "table_cell":
            slide_idx = scope.get("slide_index")
            table_index = scope.get("table_index", 0)
            row_index = scope.get("row_index", 0)
            cell_index = scope.get("cell_index", 0)
            if slide_idx is not None and slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                if table_index < len(slide.shapes):
                    table_shape = slide.shapes[table_index]
                    if getattr(table_shape, "has_table", False):
                        table = table_shape.table
                        if row_index < len(table.rows):
                            row = table.rows[row_index]
                            if cell_index < len(row.cells):
                                _set_table_cell_text(
                                    row.cells[cell_index],
                                    revised_text,
                                    font_name,
                                    font_size,
                                    align,
                                    bold,
                                    italic,
                                    underline,
                                    strike,
                                    bullet,
                                )

        elif stype == "slide":
            slide_idx = scope.get("slide_index")
            if slide_idx is not None and slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                text_shapes = [s for s in slide.shapes if s.has_text_frame]
                revised_parts = revised_text.split("\n\n")
                for i, shape in enumerate(text_shapes):
                    _set_shape_text(shape, revised_parts[i] if i < len(revised_parts) else "", font_name, font_size, align, bold, italic, underline, strike, bullet)

        elif stype == "document":
            parts = revised_text.split("\n\n---\n\n")
            for slide_idx, slide in enumerate(prs.slides):
                text_shapes = [s for s in slide.shapes if s.has_text_frame]
                slide_text = parts[slide_idx] if slide_idx < len(parts) else ""
                slide_parts = slide_text.split("\n\n")
                for i, shape in enumerate(text_shapes):
                    _set_shape_text(shape, slide_parts[i] if i < len(slide_parts) else "", font_name, font_size, align, bold, italic, underline, strike, bullet)

    prs.save(output_path)


# ── Slide structural helpers ──────────────────────────────────────────────────

def _insert_blank_slide(prs: Any, after_index: int) -> None:
    """Insert a new blank slide after *after_index*."""
    # Pick the blank layout (layout index 6 is typically 'Blank', fallback to 0)
    layout = None
    for candidate in prs.slide_layouts:
        if candidate.name.lower() in ('blank', 'blank slide', 'leer'):
            layout = candidate
            break
    if layout is None:
        layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]

    new_slide = prs.slides.add_slide(layout)

    # Move the new slide (which was appended to the end) to after_index + 1
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    new_id = slide_ids[-1]
    xml_slides.remove(new_id)
    insert_pos = min(after_index + 1, len(slide_ids) - 1)
    xml_slides.insert(insert_pos, new_id)


def _delete_slide(prs: Any, index: int) -> None:
    """Delete the slide at *index*."""
    if len(prs.slides) <= 1:
        return  # Never delete the last slide
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    if index >= len(slide_ids):
        return
    slide = prs.slides[index]
    rId = prs.slides._sldIdLst[index].get('r:id') or prs.slides._sldIdLst[index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    # Drop the relationship and remove from id list
    prs.part.drop_rel(rId)
    xml_slides.remove(slide_ids[index])


def _duplicate_slide(prs: Any, index: int) -> None:
    """Duplicate the slide at *index*, inserting the copy immediately after."""
    if index >= len(prs.slides):
        return

    # We'll copy the slide XML and add it via a new relationship
    src_slide = prs.slides[index]

    # Clone the slide XML
    cloned_xml = copy.deepcopy(src_slide._element)

    # Use the blank layout as the base for the new slide part; we'll overwrite XML
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Replace new slide's XML element children with clone
    new_sp_tree = cloned_xml.find('.//' + qn('p:spTree'))
    existing_sp_tree = new_slide._element.find('.//' + qn('p:spTree'))
    if new_sp_tree is not None and existing_sp_tree is not None:
        parent = existing_sp_tree.getparent()
        parent.replace(existing_sp_tree, copy.deepcopy(new_sp_tree))

    # Move new slide to immediately after the source
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    new_id = slide_ids[-1]
    xml_slides.remove(new_id)
    xml_slides.insert(index + 1, new_id)


# ── Text / formatting helpers ─────────────────────────────────────────────────

def _set_bullet(para: Any, enabled: bool) -> None:
    """Enable or disable bullet on a paragraph element."""
    pPr = para._p.get_or_add_pPr()
    # Remove any existing bullet elements
    for tag in (qn('a:buNone'), qn('a:buChar'), qn('a:buAutoNum'), qn('a:buClr'),
                qn('a:buFont'), qn('a:buSzPct'), qn('a:buSzPts')):
        for el in pPr.findall(tag):
            pPr.remove(el)

    if enabled:
        # Add a simple bullet character (•)
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', 'Arial')
        buFont.set('charset', '0')
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '•')
    else:
        buNone = etree.SubElement(pPr, qn('a:buNone'))


def _apply_font(run: Any, font_name: str | None, font_size: float | None,
                bold: bool | None, italic: bool | None, underline: bool | None,
                strike: bool | None) -> None:
    if font_name is not None:
        run.font.name = font_name
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if underline is not None:
        run.font.underline = underline
    if strike is not None:
        run.font.strike = strike


def _set_shape_text(shape: Any, text: str, font_name: str | None = None, font_size: float | None = None,
                    align: str | None = None, bold: bool | None = None, italic: bool | None = None,
                    underline: bool | None = None, strike: bool | None = None,
                    bullet: bool | None = None) -> None:
    _set_text_frame_text(shape.text_frame, text, font_name, font_size, align, bold, italic, underline, strike, bullet)


def _set_table_cell_text(cell: Any, text: str, font_name: str | None = None, font_size: float | None = None,
                         align: str | None = None, bold: bool | None = None, italic: bool | None = None,
                         underline: bool | None = None, strike: bool | None = None,
                         bullet: bool | None = None) -> None:
    _set_text_frame_text(cell.text_frame, text, font_name, font_size, align, bold, italic, underline, strike, bullet)


def _set_text_frame_text(tf: Any, text: str, font_name: str | None = None, font_size: float | None = None,
                         align: str | None = None, bold: bool | None = None, italic: bool | None = None,
                         underline: bool | None = None, strike: bool | None = None,
                         bullet: bool | None = None) -> None:
    if not tf.paragraphs:
        return
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = text
        for run in first_para.runs[1:]:
            run.text = ""
        runs = first_para.runs
    else:
        runs = [first_para.add_run()]
        runs[0].text = text

    if align is not None and align in _ALIGN_MAP:
        first_para.alignment = _ALIGN_MAP[align]

    if bullet is not None:
        _set_bullet(first_para, bullet)

    for run in runs:
        _apply_font(run, font_name, font_size, bold, italic, underline, strike)
    # Remove extra paragraphs
    for para in tf.paragraphs[1:]:
        p_elem = para._p
        p_elem.getparent().remove(p_elem)
